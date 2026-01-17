#!/usr/bin/env python3
"""
Convert PowerPoint files to Markdown + Images for LLM consumption.

Usage:
    python pptx_to_markdown.py <input.pptx> [output_dir]

Output:
    - output_dir/slides/slide-01.png, slide-02.png, ...
    - output_dir/content.md (markdown with text + image references)
"""

import argparse
import os
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_text_from_slide(slide):
    """Extract all text content from a slide."""
    texts = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # Determine indent level
                    level = paragraph.level if hasattr(paragraph, 'level') else 0
                    prefix = "  " * level + "- " if level > 0 else ""
                    texts.append(f"{prefix}{text}")

        # Handle tables
        if shape.has_table:
            table = shape.table
            table_md = extract_table_as_markdown(table)
            if table_md:
                texts.append(table_md)

    return texts


def extract_table_as_markdown(table):
    """Convert a PowerPoint table to markdown format."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.strip().replace("|", "\\|")
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    # Build markdown table
    md_lines = []
    md_lines.append("| " + " | ".join(rows[0]) + " |")
    md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    for row in rows[1:]:
        md_lines.append("| " + " | ".join(row) + " |")

    return "\n".join(md_lines)


def extract_speaker_notes(slide):
    """Extract speaker notes from a slide."""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text.strip()
        if notes_text:
            return notes_text
    return None


def export_slides_as_images(pptx_path, output_dir):
    """Use LibreOffice to export slides as PNG images."""
    slides_dir = output_dir / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    # Create a temporary directory for LibreOffice output
    with tempfile.TemporaryDirectory() as temp_dir:
        # Convert PPTX to PDF first (better quality)
        subprocess.run([
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", temp_dir,
            str(pptx_path)
        ], check=True, capture_output=True)

        # Find the PDF
        pdf_name = pptx_path.stem + ".pdf"
        pdf_path = Path(temp_dir) / pdf_name

        if not pdf_path.exists():
            print(f"Warning: PDF conversion failed, trying direct PNG export")
            # Fallback: try direct PNG export
            subprocess.run([
                "soffice",
                "--headless",
                "--convert-to", "png",
                "--outdir", str(slides_dir),
                str(pptx_path)
            ], check=True, capture_output=True)
            return

        # Convert PDF pages to PNG using sips (macOS) or ImageMagick
        # First, let's try using pdftoppm if available, otherwise use sips
        try:
            # Try pdftoppm (from poppler)
            subprocess.run([
                "pdftoppm",
                "-png",
                "-r", "150",  # 150 DPI for good quality without huge files
                str(pdf_path),
                str(slides_dir / "slide")
            ], check=True, capture_output=True)

            # Rename files to have consistent naming (slide-01.png format)
            for f in slides_dir.glob("slide-*.png"):
                # pdftoppm creates slide-1.png, slide-2.png, etc.
                # We want slide-01.png, slide-02.png
                parts = f.stem.split("-")
                if len(parts) == 2:
                    num = int(parts[1])
                    new_name = f"slide-{num:02d}.png"
                    f.rename(slides_dir / new_name)

        except (subprocess.CalledProcessError, FileNotFoundError):
            # Fallback: use ImageMagick's convert
            try:
                subprocess.run([
                    "convert",
                    "-density", "150",
                    str(pdf_path),
                    str(slides_dir / "slide-%02d.png")
                ], check=True, capture_output=True)
            except (subprocess.CalledProcessError, FileNotFoundError):
                print("Warning: Could not convert PDF to images. Install poppler or imagemagick.")
                print("  brew install poppler")
                print("  or: brew install imagemagick")
                # Copy the PDF as fallback
                shutil.copy(pdf_path, output_dir / "slides.pdf")


def convert_pptx_to_markdown(pptx_path, output_dir):
    """Main conversion function."""
    pptx_path = Path(pptx_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Load the presentation
    prs = Presentation(str(pptx_path))

    # Export slides as images
    print(f"Exporting slides as images...")
    export_slides_as_images(pptx_path, output_dir)

    # Build markdown content
    md_lines = []
    md_lines.append(f"# {pptx_path.stem}")
    md_lines.append("")
    md_lines.append(f"*Converted from: {pptx_path.name}*")
    md_lines.append("")
    md_lines.append("---")
    md_lines.append("")

    for i, slide in enumerate(prs.slides, 1):
        slide_num = f"{i:02d}"

        # Slide header
        md_lines.append(f"## Slide {i}")
        md_lines.append("")

        # Image reference
        md_lines.append(f"![Slide {i}](./slides/slide-{slide_num}.png)")
        md_lines.append("")

        # Extract text content
        texts = extract_text_from_slide(slide)
        if texts:
            md_lines.append("### Content")
            md_lines.append("")
            for text in texts:
                md_lines.append(text)
            md_lines.append("")

        # Extract speaker notes
        notes = extract_speaker_notes(slide)
        if notes:
            md_lines.append("### Speaker Notes")
            md_lines.append("")
            md_lines.append(f"> {notes}")
            md_lines.append("")

        md_lines.append("---")
        md_lines.append("")

    # Write markdown file
    md_path = output_dir / "content.md"
    with open(md_path, "w", encoding="utf-8") as f:
        f.write("\n".join(md_lines))

    print(f"Done! Output written to: {output_dir}")
    print(f"  - Markdown: {md_path}")
    print(f"  - Images: {output_dir / 'slides'}/")

    return md_path


def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint to Markdown + Images for LLM consumption"
    )
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument(
        "output",
        nargs="?",
        help="Output directory (default: same name as input file)"
    )

    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: File not found: {input_path}")
        return 1

    if args.output:
        output_dir = Path(args.output)
    else:
        output_dir = input_path.parent / input_path.stem

    convert_pptx_to_markdown(input_path, output_dir)
    return 0


if __name__ == "__main__":
    exit(main())
